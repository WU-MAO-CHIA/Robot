from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# 建立簡報
prs = Presentation()

# --- 第一頁：標題 + 條列 ---
slide_layout = prs.slide_layouts[1]  # Title and Content
slide = prs.slides.add_slide(slide_layout)

# 標題
title = slide.shapes.title
title.text = "積木化測試腳本，降低維護成本"

# 條列內容
content = slide.placeholders[1].text_frame
content.clear()
points = [
    "模組化設計：將測試腳本拆解為可重複使用的積木模組",
    "快速拼裝：依測試需求靈活組合，不需重複撰寫",
    "降低維護成本：修改集中於模組，減少整體維護工作量",
    "提升效率：縮短腳本更新時間，加快測試迭代"
]
for point in points:
    p = content.add_paragraph()
    p.text = point
    p.level = 0

# --- 第二頁：積木示意圖 ---
blank_slide_layout = prs.slide_layouts[6]  # Blank
slide2 = prs.slides.add_slide(blank_slide_layout)

# 標題
txBox = slide2.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
tf = txBox.text_frame
p = tf.add_paragraph()
p.text = "積木化測試腳本示意圖"
p.font.size = Pt(28)
p.font.bold = True
p.alignment = PP_ALIGN.CENTER

# 彩色矩形模擬積木
colors = [RGBColor(255, 99, 71), RGBColor(135, 206, 250), RGBColor(144, 238, 144), RGBColor(255, 215, 0)]
labels = ["登入模組", "交易模組", "查詢模組", "報表模組"]

for i, (color, label) in enumerate(zip(colors, labels)):
    left = Inches(1.5 + i*1.8)
    top = Inches(2.5)
    width = Inches(1.6)
    height = Inches(1)
    shape = slide2.shapes.add_shape(1, left, top, width, height)  # 1=矩形
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.text = label
    shape.text_frame.paragraphs[0].font.size = Pt(14)
    shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# 儲存
prs.save("積木化測試腳本說明_含示意圖.pptx")
print("已輸出檔案：積木化測試腳本說明_含示意圖.pptx")